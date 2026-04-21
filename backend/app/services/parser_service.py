from pathlib import Path

import fitz
from pptx import Presentation

from app.services.storage_service import LocalStorageService


class ParserService:
    def __init__(self) -> None:
        self.storage = LocalStorageService()

    def extract_text(self, storage_key: str, content_type: str) -> str:
        resolved_path = self.storage.resolve_storage_path(storage_key)
        suffix = resolved_path.suffix.lower()
        if content_type == "application/pdf" or suffix == ".pdf":
            return self._extract_pdf(resolved_path)
        if suffix in {".ppt", ".pptx"} or "presentation" in content_type:
            return self._extract_pptx(resolved_path)
        return resolved_path.read_text(encoding="utf-8", errors="ignore")

    def _extract_pdf(self, storage_path: Path) -> str:
        document = fitz.open(storage_path)
        pages = [page.get_text("text") for page in document]
        document.close()
        return "\n".join(text.strip() for text in pages if text.strip())

    def _extract_pptx(self, storage_path: Path) -> str:
        presentation = Presentation(storage_path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = [f"Slide {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text.strip())
            notes_frame = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
            if notes_frame and notes_frame.text:
                parts.append(f"Notes: {notes_frame.text.strip()}")
            slides.append("\n".join(part for part in parts if part))
        return "\n\n".join(slides)
